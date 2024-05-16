from fastapi import FastAPI, Response
from pydantic import BaseModel
from openai import OpenAI
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import ast

format = {
    "slide title": "<Title of the Presentation>",
    "<Slide 1 name>": [
        "<Point 1 for Slide 1>",
        "<Point 2 for Slide 1>",
        # Add more slides as needed
    ],
    "<Slide 2 name>": [
        "<Point 1 for Slide 2>",
        "<Point 2 for Slide 2>",
        "<Point 3 for Slide 2>",
        "<Point 4 for Slide 2>",
        # Add more slides as needed
    ],
    "<Slide 3 name>": [
        "<Point 1 for Slide 3>",
        "<Point 2 for Slide 3>",
        "<Point 3 for Slide 3>",
        # Add more slides as needed
    ],
    # Add more slides as needed
}


app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


class Item(BaseModel):
    topic: str
    num_slides: int


@app.post("/generate_presentation/")
async def generate_presentation(item: Item):
    api_key = "testapikey"
    client = OpenAI(api_key=api_key)
    presentation_content = client.chat.completions.create(
        model="gpt-3.5-turbo-0125",
        messages=[
            {"role": "system", "content": "You are a helpful presentation assistant."},
            {
                "role": "user",
                "content": f"Create a {item.num_slides} slide presentation on {item.topic} strictly in the following format : {format}, No need to mention Slide 1, slide 2 etc , just direct names",
            },
        ],
    )
    slides_content = presentation_content.choices[0].message.content

    prs = Presentation()

    presentation_content_dict = ast.literal_eval(slides_content)
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = presentation_content_dict["slide title"]

    for title, content in list(presentation_content_dict.items())[1:]:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        for line in content:
            p = slide.placeholders[1].text_frame.add_paragraph()
            p.text = line

    ppt_io = BytesIO()
    prs.save(ppt_io)

    return Response(
        content=ppt_io.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    # return 1
